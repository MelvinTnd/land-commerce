"""
BeninMarket - Présentation PowerPoint Premium
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette de couleurs ───────────────────────────
VERT       = RGBColor(0x1B, 0x6B, 0x3A)  # Vert BeninMarket
VERT_DARK  = RGBColor(0x0D, 0x4A, 0x28)  # Vert foncé
VERT_LIGHT = RGBColor(0xE6, 0xF8, 0xEA)  # Vert clair
OR         = RGBColor(0xD4, 0x92, 0x0A)  # Doré
OR_LIGHT   = RGBColor(0xFE, 0xF3, 0xC7)  # Doré clair
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
NOIR       = RGBColor(0x0D, 0x0D, 0x0D)
GRIS       = RGBColor(0x9C, 0xA3, 0xAF)
GRIS_BG    = RGBColor(0xF7, 0xF5, 0xF0)
GRIS_CARD  = RGBColor(0xF3, 0xF4, 0xF6)
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)
ROSE       = RGBColor(0xDB, 0x27, 0x77)

W = Inches(13.333)   # Largeur slide 16:9
H = Inches(7.5)      # Hauteur slide 16:9

def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=NOIR,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def new_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Virer les placeholders par défaut
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide

# ════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)

# Fond vert dégradé (2 rectangles)
add_rect(slide, 0, 0, W, H, VERT_DARK)
add_rect(slide, 0, 0, W, Inches(5), VERT)

# Motif cercles décoratifs
for cx, cy, sz, op in [(12.5, 0.5, 3.5, 0.07), (0, 5.5, 2.5, 0.05), (7, 6, 4, 0.04)]:
    s = slide.shapes.add_shape(9, Inches(cx)-Inches(sz/2), Inches(cy)-Inches(sz/2),
                                Inches(sz), Inches(sz))
    s.fill.solid(); s.fill.fore_color.rgb = BLANC
    s.line.fill.background()
    from pptx.oxml.ns import qn
    from lxml import etree
    α = int(op * 100000)
    s.fill.fore_color._color._xClr.set('lastClr', 'FFFFFF')

# Bande dorée accents
add_rect(slide, 0, Inches(5.0), W, Inches(0.12), OR)

# Logo texte
add_text(slide, "🌿 BéninMarket", Inches(0.8), Inches(1.2), Inches(7), Inches(1.2),
         size=54, bold=True, color=BLANC, font="Calibri")

# Tagline
add_text(slide, "La Marketplace Locale du Bénin",
         Inches(0.85), Inches(2.5), Inches(9), Inches(0.7),
         size=24, bold=False, color=RGBColor(0xBB,0xF7,0xD0), font="Calibri")

# Badge
badge = add_rect(slide, Inches(0.85), Inches(3.3), Inches(3.5), Inches(0.5), OR)
badge.line.fill.background()
add_text(slide, "Projet Full-Stack 2026",
         Inches(0.9), Inches(3.33), Inches(3.4), Inches(0.45),
         size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER, font="Calibri")

# Infos bas de page
add_text(slide, "Next.js  ·  Laravel 11  ·  MySQL  ·  Render  ·  Vercel",
         Inches(0.85), Inches(5.3), Inches(9), Inches(0.5),
         size=14, color=RGBColor(0x86,0xEF,0xAC), font="Calibri")

add_text(slide, "Présenté par : Melvin TONADO   |   Avril 2026",
         Inches(0.85), Inches(6.6), Inches(9), Inches(0.5),
         size=12, color=RGBColor(0x9C,0xA3,0xAF), font="Calibri")

# Numéro slide
add_text(slide, "1", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, GRIS_BG)
add_rect(slide, 0, 0, Inches(0.35), H, VERT)

add_text(slide, "Table des Matières", Inches(0.65), Inches(0.5), Inches(10), Inches(1),
         size=36, bold=True, color=NOIR, font="Calibri")
add_rect(slide, Inches(0.65), Inches(1.35), Inches(2.5), Inches(0.06), OR)

items = [
    ("01", "Contexte & Problématique",     "Pourquoi BeninMarket ?"),
    ("02", "Solution & Vision",             "Notre réponse au marché"),
    ("03", "Fonctionnalités Clés",          "Ce que la plateforme offre"),
    ("04", "Architecture Technique",        "Stack & infrastructure"),
    ("05", "Interface Utilisateur",         "Design & expérience"),
    ("06", "Dashboard Vendeur",             "Espace de gestion"),
    ("07", "API Backend",                   "Routes & sécurité"),
    ("08", "Déploiement",                   "Render + Vercel"),
    ("09", "Perspectives",                  "Évolutions futures"),
]
for i, (num, titre, sub) in enumerate(items):
    col = i % 3
    row = i // 3
    bx = Inches(0.65 + col * 4.2)
    by = Inches(1.6 + row * 1.8)
    card = add_rect(slide, bx, by, Inches(3.9), Inches(1.5), BLANC)
    # Accent couleur
    c = [VERT, OR, VIOLET][col]
    add_rect(slide, bx, by, Inches(0.07), Inches(1.5), c)
    add_text(slide, num, bx+Inches(0.2), by+Inches(0.1), Inches(0.7), Inches(0.45),
             size=22, bold=True, color=c, font="Calibri")
    add_text(slide, titre, bx+Inches(0.2), by+Inches(0.5), Inches(3.5), Inches(0.45),
             size=12, bold=True, color=NOIR, font="Calibri")
    add_text(slide, sub, bx+Inches(0.2), by+Inches(0.95), Inches(3.5), Inches(0.45),
             size=10, color=GRIS, font="Calibri")

add_text(slide, "2", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE & PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, GRIS_BG)
add_rect(slide, 0, 0, Inches(0.35), H, VERT)

# Header
add_rect(slide, Inches(0.65), Inches(0.5), Inches(1), Inches(1), VERT_LIGHT)
add_text(slide, "01", Inches(0.65), Inches(0.5), Inches(1), Inches(1),
         size=28, bold=True, color=VERT, align=PP_ALIGN.CENTER, font="Calibri")
add_text(slide, "Contexte & Problématique",
         Inches(1.9), Inches(0.55), Inches(9), Inches(0.9),
         size=32, bold=True, color=NOIR, font="Calibri")
add_rect(slide, Inches(1.9), Inches(1.4), Inches(4), Inches(0.06), VERT)

# Problèmes
problems = [
    ("🏪", "Marché fragmenté", "Les artisans et commerçants béninois n'ont pas de vitrine numérique unifiée pour toucher leurs clients."),
    ("💸", "Paiements limités", "L'absence de solutions de paiement mobile locales (MTN, Moov) freine l'essor du commerce en ligne."),
    ("📦", "Gestion manuelle", "Les vendeurs gèrent stocks, commandes et clients de façon artisanale sans outils digitaux adaptés."),
    ("🌐", "Visibilité faible", "Les boutiques locales souffrent d'un manque de visibilité face aux grandes plateformes internationales."),
]
for i, (ico, titre, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    bx = Inches(0.65 + col * 6.3)
    by = Inches(1.7 + row * 2.5)
    add_rect(slide, bx, by, Inches(6), Inches(2.2), BLANC)
    add_text(slide, ico, bx+Inches(0.25), by+Inches(0.25), Inches(0.7), Inches(0.7), size=28)
    add_text(slide, titre, bx+Inches(1.1), by+Inches(0.2), Inches(4.6), Inches(0.5),
             size=14, bold=True, color=NOIR, font="Calibri")
    add_rect(slide, bx+Inches(1.1), by+Inches(0.7), Inches(0.6), Inches(0.04), OR)
    add_text(slide, desc, bx+Inches(0.25), by+Inches(0.9), Inches(5.5), Inches(1.1),
             size=11, color=GRIS, font="Calibri")

add_text(slide, "3", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION & VISION
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
# Split layout : gauche vert, droite clair
add_rect(slide, 0, 0, Inches(6), H, VERT)
add_rect(slide, Inches(6), 0, Inches(7.333), H, GRIS_BG)

# Gauche
add_text(slide, "02", Inches(0.5), Inches(0.5), Inches(1.5), Inches(1),
         size=48, bold=True, color=RGBColor(0xBB,0xF7,0xD0), font="Calibri")
add_text(slide, "Notre Solution",
         Inches(0.5), Inches(1.4), Inches(5), Inches(0.8),
         size=34, bold=True, color=BLANC, font="Calibri")
add_rect(slide, Inches(0.5), Inches(2.15), Inches(3), Inches(0.08), OR)
add_text(slide,
    "BéninMarket est une marketplace multi-vendeurs 100 % locale, "
    "conçue pour connecter artisans, commerçants et producteurs béninois "
    "à leurs clients, en offrant une expérience premium adaptée au contexte africain.",
    Inches(0.5), Inches(2.4), Inches(5.2), Inches(2.5),
    size=13, color=RGBColor(0xBB,0xF7,0xD0), font="Calibri")

# Stats clés
stats = [("3 000+", "Utilisateurs cibles"), ("500+", "Produits référencés"), ("100%", "Made in Bénin")]
for i, (val, lbl) in enumerate(stats):
    by = Inches(5.0 + i * 0.7)
    add_text(slide, val + "  " + lbl, Inches(0.5), by, Inches(5), Inches(0.6),
             size=14, bold=True, color=OR, font="Calibri")

# Droite — piliers
add_text(slide, "Nos Piliers", Inches(6.4), Inches(0.6), Inches(6.5), Inches(0.8),
         size=28, bold=True, color=NOIR, font="Calibri")

piliers = [
    (VERT,   "accessibility", "Accessibilité",  "Interface simple, optimisée mobile & desktop"),
    (OR,     "payments",      "Paiement Local",  "MTN MoMo & Moov Money intégrés"),
    (VIOLET, "storefront",    "Multi-Vendeurs",  "Chaque vendeur gère sa boutique en autonomie"),
    (ROSE,   "analytics",     "Tableau de bord", "Statistiques & gestion complète du catalogue"),
]
for i, (c, icon, titre, desc) in enumerate(piliers):
    by = Inches(1.5 + i * 1.35)
    add_rect(slide, Inches(6.4), by, Inches(6.6), Inches(1.15), BLANC)
    add_rect(slide, Inches(6.4), by, Inches(0.08), Inches(1.15), c)
    icbox = add_rect(slide, Inches(6.7), by+Inches(0.3), Inches(0.55), Inches(0.55), RGBColor(c.red//5*4, c.green//5*4, c.blue//5*4) if True else c)
    icbox.fill.solid(); icbox.fill.fore_color.rgb = RGBColor(int(c.red*0.15+0xE6*0.85), int(c.green*0.15+0xE6*0.85), int(c.blue*0.15+0xE6*0.85))
    add_text(slide, titre, Inches(7.45), by+Inches(0.15), Inches(5.2), Inches(0.45),
             size=13, bold=True, color=NOIR, font="Calibri")
    add_text(slide, desc, Inches(7.45), by+Inches(0.6), Inches(5.2), Inches(0.45),
             size=11, color=GRIS, font="Calibri")

add_text(slide, "4", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — FONCTIONNALITÉS CLÉS
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, GRIS_BG)
add_rect(slide, 0, 0, Inches(0.35), H, OR)

add_text(slide, "Fonctionnalités Clés", Inches(0.65), Inches(0.45), Inches(10), Inches(0.9),
         size=32, bold=True, color=NOIR, font="Calibri")
add_rect(slide, Inches(0.65), Inches(1.3), Inches(3.5), Inches(0.07), OR)
add_text(slide, "03", Inches(11.8), Inches(0.45), Inches(1.2), Inches(0.8),
         size=36, bold=True, color=OR, align=PP_ALIGN.RIGHT, font="Calibri")

features = [
    ("🛍️", "Marketplace",         "Catalogue de produits filtrable par catégorie, prix et recherche textuelle en temps réel."),
    ("🏪", "Boutiques",            "Pages boutiques personnalisées par vendeur avec produits, avis et informations."),
    ("🎁", "Promotions",           "Section ventes flash avec compte à rebours, gérée depuis le dashboard admin."),
    ("📦", "Inventaire",           "Ajout, suppression et gestion du stock produits directement depuis l'espace vendeur."),
    ("📊", "Dashboard Vendeur",    "Statistiques, commandes, paramètres boutique et solde dans un espace dédié."),
    ("📝", "Blog & Forum",         "Articles et sujets de discussion pour la communauté de vendeurs et acheteurs."),
    ("🔐", "Authentification",     "Inscription, connexion, rôles (acheteur / vendeur) via JWT Sanctum sécurisé."),
    ("🌙", "Design Premium",       "Interface moderne dark/light mode, animations fluides, 100% responsive."),
    ("🚀", "Performance",          "Rendu hybride SSR/Client, chargement optimisé, API REST rapide sur Render."),
]
cols = 3
per_col = 3
for idx, (ico, titre, desc) in enumerate(features):
    col = idx % cols
    row = idx // cols
    bx = Inches(0.55 + col * 4.25)
    by = Inches(1.55 + row * 1.9)
    card = add_rect(slide, bx, by, Inches(4.0), Inches(1.7), BLANC)
    ico_bg = [VERT_LIGHT, OR_LIGHT, RGBColor(0xED,0xE9,0xFE)][col]
    add_rect(slide, bx+Inches(0.2), by+Inches(0.2), Inches(0.6), Inches(0.6), ico_bg)
    add_text(slide, ico, bx+Inches(0.22), by+Inches(0.18), Inches(0.6), Inches(0.6), size=18)
    add_text(slide, titre, bx+Inches(1.0), by+Inches(0.15), Inches(2.8), Inches(0.45),
             size=12, bold=True, color=NOIR, font="Calibri")
    add_text(slide, desc, bx+Inches(0.2), by+Inches(0.9), Inches(3.6), Inches(0.65),
             size=9.5, color=GRIS, font="Calibri")

add_text(slide, "5", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE TECHNIQUE
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, NOIR)
add_rect(slide, 0, 0, W, Inches(0.08), VERT)
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), OR)

add_text(slide, "04 — Architecture Technique",
         Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         size=30, bold=True, color=BLANC, font="Calibri")

# Stack layers
layers = [
    ("FRONTEND",   "Next.js 14  ·  Tailwind CSS  ·  NextAuth.js",      VERT,   "Vercel CDN"),
    ("API REST",   "Laravel 11  ·  Sanctum JWT  ·  Eloquent ORM",       OR,     "Render.com"),
    ("BASE DE DONNÉES", "MySQL (Railway)  ·  Migrations  ·  Seeds",     VIOLET, "Railway.app"),
    ("MÉDIAS",     "Images Unsplash  ·  Upload local  ·  Optimisation", ROSE,   "Local Storage"),
]
for i, (lbl, tech, color, host) in enumerate(layers):
    by = Inches(1.3 + i * 1.4)
    add_rect(slide, Inches(0.6), by, Inches(1.5), Inches(1.15), color)
    add_text(slide, lbl, Inches(0.6), by+Inches(0.35), Inches(1.5), Inches(0.5),
             size=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER, font="Calibri")
    add_rect(slide, Inches(2.15), by, Inches(7.8), Inches(1.15), RGBColor(0x1A,0x1A,0x2E))
    add_text(slide, tech, Inches(2.4), by+Inches(0.35), Inches(6.5), Inches(0.5),
             size=14, bold=True, color=BLANC, font="Calibri")
    # Badge host
    add_rect(slide, Inches(10.2), by+Inches(0.3), Inches(2.6), Inches(0.5), color)
    add_text(slide, "☁ " + host, Inches(10.2), by+Inches(0.3), Inches(2.6), Inches(0.5),
             size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER, font="Calibri")

# Flèches de connexion (traits)
for i in range(3):
    by = Inches(1.3 + i * 1.4) + Inches(1.15)
    line = slide.shapes.add_connector(1, Inches(1.35), by, Inches(1.35), by + Inches(0.25))
    line.line.color.rgb = VERT
    line.line.width = Pt(2)

add_text(slide, "6", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — DASHBOARD VENDEUR
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, GRIS_BG)
add_rect(slide, 0, 0, Inches(0.35), H, VIOLET)

add_text(slide, "Dashboard Vendeur", Inches(0.65), Inches(0.45), Inches(10), Inches(0.9),
         size=32, bold=True, color=NOIR, font="Calibri")
add_rect(slide, Inches(0.65), Inches(1.3), Inches(3.5), Inches(0.07), VIOLET)
add_text(slide, "06", Inches(11.8), Inches(0.45), Inches(1.2), Inches(0.8),
         size=36, bold=True, color=VIOLET, align=PP_ALIGN.RIGHT, font="Calibri")

# KPI cards
kpis = [
    ("Revenus", "0 CFA", "+12%", VERT, VERT_LIGHT),
    ("Produits", "4 articles", "", VIOLET, RGBColor(0xED,0xE9,0xFE)),
    ("Commandes", "0 reçues", "", OR, OR_LIGHT),
    ("Note moy.", "4.9 / 5", "Top 1%", ROSE, RGBColor(0xFC,0xE7,0xF3)),
]
for i, (lbl, val, badge, c, bg) in enumerate(kpis):
    bx = Inches(0.65 + i * 3.15)
    add_rect(slide, bx, Inches(1.55), Inches(2.9), Inches(1.3), BLANC)
    add_rect(slide, bx+Inches(0.25), Inches(1.75), Inches(0.6), Inches(0.6), bg)
    add_text(slide, lbl, bx+Inches(0.25), Inches(2.45), Inches(2.4), Inches(0.35),
             size=9, bold=True, color=GRIS, font="Calibri")
    add_text(slide, val, bx+Inches(0.25), Inches(2.75), Inches(2.4), Inches(0.4),
             size=14, bold=True, color=NOIR, font="Calibri")
    if badge:
        add_rect(slide, bx+Inches(2.0), Inches(1.7), Inches(0.65), Inches(0.3), bg)
        add_text(slide, badge, bx+Inches(2.0), Inches(1.7), Inches(0.65), Inches(0.3),
                 size=8, bold=True, color=c, align=PP_ALIGN.CENTER, font="Calibri")

# Onglets
tabs = [
    ("📊", "Tableau de bord",     "Vue globale des performances et activité récente"),
    ("📦", "Inventaire & Stock",   "Ajouter, modifier et supprimer des produits"),
    ("🧾", "Commandes",            "Historique et suivi des commandes clients"),
    ("⚙️", "Paramètres",           "Modifier nom, logo, description de la boutique"),
]
for i, (ico, titre, desc) in enumerate(tabs):
    bx = Inches(0.65 + (i%2)*6.3)
    by = Inches(3.1 + (i//2)*1.6)
    add_rect(slide, bx, by, Inches(6.0), Inches(1.35), BLANC)
    add_rect(slide, bx, by, Inches(0.07), Inches(1.35), VIOLET)
    add_text(slide, ico + "  " + titre, bx+Inches(0.25), by+Inches(0.2), Inches(5.5), Inches(0.5),
             size=13, bold=True, color=NOIR, font="Calibri")
    add_text(slide, desc, bx+Inches(0.25), by+Inches(0.72), Inches(5.5), Inches(0.45),
             size=11, color=GRIS, font="Calibri")

add_text(slide, "7", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — API BACKEND
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, RGBColor(0x0F,0x1A,0x0F))
add_rect(slide, 0, 0, W, Inches(0.06), VERT)

add_text(slide, "07 — API Backend", Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         size=30, bold=True, color=BLANC, font="Calibri")
add_text(slide, "Laravel 11  ·  Sanctum JWT  ·  REST API",
         Inches(0.6), Inches(1.0), Inches(10), Inches(0.5),
         size=14, color=RGBColor(0x86,0xEF,0xAC), font="Calibri")

endpoints = [
    ("PUBLIC",            "GET  /api/products", "GET  /api/shops", "GET  /api/promotions", "GET  /api/categories"),
    ("AUTH",              "POST /api/register", "POST /api/login", "POST /api/logout", "GET  /api/user"),
    ("VENDEUR 🔐",        "GET  /api/vendor/products", "POST /api/vendor/products", "DELETE /api/vendor/products/{id}", "GET  /api/vendor/dashboard"),
    ("COMMANDES 🔐",      "POST /api/checkout", "GET  /api/orders", "GET  /api/orders/{id}", "POST /api/orders/{id}/cancel"),
]
for i, (titre, *routes) in enumerate(endpoints):
    col = i % 2
    row = i // 2
    bx = Inches(0.5 + col * 6.4)
    by = Inches(1.7 + row * 2.6)
    add_rect(slide, bx, by, Inches(6.1), Inches(2.35), RGBColor(0x1A,0x2E,0x1A))
    c = [VERT, OR, VIOLET, ROSE][i]
    add_rect(slide, bx, by, Inches(0.08), Inches(2.35), c)
    add_text(slide, titre, bx+Inches(0.2), by+Inches(0.1), Inches(5.6), Inches(0.45),
             size=12, bold=True, color=c, font="Calibri")
    for j, route in enumerate(routes):
        add_text(slide, route, bx+Inches(0.2), by+Inches(0.55+j*0.42), Inches(5.6), Inches(0.38),
                 size=10.5, color=RGBColor(0xBB,0xF7,0xD0), font="Courier New")

add_text(slide, "8", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — DÉPLOIEMENT
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, GRIS_BG)
add_rect(slide, 0, 0, Inches(0.35), H, OR)

add_text(slide, "Déploiement", Inches(0.65), Inches(0.45), Inches(10), Inches(0.9),
         size=32, bold=True, color=NOIR, font="Calibri")
add_rect(slide, Inches(0.65), Inches(1.3), Inches(2), Inches(0.07), OR)
add_text(slide, "08", Inches(11.8), Inches(0.45), Inches(1.2), Inches(0.8),
         size=36, bold=True, color=OR, align=PP_ALIGN.RIGHT, font="Calibri")

deploy = [
    ("☁", "Vercel – Frontend",
     "land-commerce.vercel.app",
     ["Déploiement automatique depuis GitHub", "CDN mondial edge", "Environnement Next.js optimisé", "Variables d'environnement sécurisées"],
     VERT),
    ("🚀", "Render – Backend",
     "land-commerce-api.onrender.com",
     ["Docker containerisé (PHP 8.2 + Apache)", "Migrations auto au démarrage", "Base MySQL Railway", "CI/CD GitHub Actions intégré"],
     VIOLET),
    ("🗄️", "Railway – Base de données",
     "metro.proxy.rlwy.net",
     ["MySQL 8.0 managé", "Backups automatiques", "SSL activé", "Accès par proxy sécurisé"],
     OR),
]
for i, (ico, titre, url, items, c) in enumerate(deploy):
    bx = Inches(0.55 + i * 4.25)
    add_rect(slide, bx, Inches(1.55), Inches(4.0), Inches(5.5), BLANC)
    add_rect(slide, bx, Inches(1.55), Inches(4.0), Inches(0.07), c)
    add_text(slide, ico + " " + titre, bx+Inches(0.25), Inches(1.75), Inches(3.5), Inches(0.55),
             size=14, bold=True, color=NOIR, font="Calibri")
    add_rect(slide, bx+Inches(0.25), Inches(2.4), Inches(3.5), Inches(0.35), GRIS_CARD)
    add_text(slide, url, bx+Inches(0.3), Inches(2.42), Inches(3.4), Inches(0.32),
             size=9, color=c, font="Courier New")
    for j, it in enumerate(items):
        add_text(slide, "✓  " + it, bx+Inches(0.25), Inches(2.95+j*0.85), Inches(3.5), Inches(0.7),
                 size=11, color=GRIS, font="Calibri")

add_text(slide, "9", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — PERSPECTIVES & ÉVOLUTIONS
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, GRIS_BG)
add_rect(slide, 0, 0, Inches(0.35), H, ROSE)

add_text(slide, "Perspectives & Évolutions", Inches(0.65), Inches(0.45), Inches(10), Inches(0.9),
         size=32, bold=True, color=NOIR, font="Calibri")
add_rect(slide, Inches(0.65), Inches(1.3), Inches(4), Inches(0.07), ROSE)
add_text(slide, "09", Inches(11.8), Inches(0.45), Inches(1.2), Inches(0.8),
         size=36, bold=True, color=ROSE, align=PP_ALIGN.RIGHT, font="Calibri")

roadmap = [
    ("Q3 2026", [
        ("💳 Paiement MTN / Moov", VERT),
        ("📱 Application Mobile React Native", OR),
    ]),
    ("Q4 2026", [
        ("🤖 Recommandations IA", VIOLET),
        ("⭐ Système d'avis avancé", ROSE),
    ]),
    ("2027", [
        ("🌍 Expansion UEMOA", VERT),
        ("📊 Analytics avancées", OR),
    ]),
]
for i, (period, items) in enumerate(roadmap):
    bx = Inches(0.65 + i * 4.2)
    add_rect(slide, bx, Inches(1.6), Inches(3.9), Inches(5.2), BLANC)
    add_rect(slide, bx, Inches(1.6), Inches(3.9), Inches(0.6), NOIR)
    add_text(slide, period, bx, Inches(1.6), Inches(3.9), Inches(0.6),
             size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER, font="Calibri")
    for j, (txt, c) in enumerate(items):
        by = Inches(2.4 + j * 2.0)
        add_rect(slide, bx+Inches(0.25), by, Inches(3.4), Inches(1.7), GRIS_CARD)
        add_rect(slide, bx+Inches(0.25), by, Inches(0.07), Inches(1.7), c)
        add_text(slide, txt, bx+Inches(0.5), by+Inches(0.6), Inches(3.0), Inches(0.6),
                 size=12, bold=True, color=NOIR, font="Calibri")

add_text(slide, "10", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSION
# ══════════════════════════════════════════════════════════════════
slide = new_slide(prs)
add_rect(slide, 0, 0, W, H, VERT_DARK)
add_rect(slide, 0, 0, W, Inches(0.1), OR)
add_rect(slide, 0, H-Inches(0.1), W, Inches(0.1), OR)

# Cercles déco
for cx, cy, sz in [(11, 1.5, 4), (1, 6, 3), (7, 4, 6)]:
    s = slide.shapes.add_shape(9, Inches(cx)-Inches(sz/2), Inches(cy)-Inches(sz/2),
                                Inches(sz), Inches(sz))
    s.fill.solid(); s.fill.fore_color.rgb = VERT
    s.line.fill.background()

add_text(slide, "🌿", Inches(5.8), Inches(0.8), Inches(1.5), Inches(1.2), size=48, align=PP_ALIGN.CENTER)
add_text(slide, "Merci de votre attention !",
         Inches(1), Inches(1.9), Inches(11.3), Inches(1.2),
         size=46, bold=True, color=BLANC, align=PP_ALIGN.CENTER, font="Calibri")
add_rect(slide, Inches(4.5), Inches(3.05), Inches(4.3), Inches(0.08), OR)
add_text(slide, "BéninMarket — La Marketplace Locale du Bénin",
         Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.7),
         size=18, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.CENTER, font="Calibri")

# Infos contact
contacts = [
    ("🌐 Frontend", "land-commerce.vercel.app", VERT),
    ("⚡ Backend", "land-commerce-api.onrender.com/api", OR),
    ("💻 GitHub", "github.com/MelvinTnd", VIOLET),
]
for i, (lbl, val, c) in enumerate(contacts):
    bx = Inches(1.0 + i * 3.8)
    add_rect(slide, bx, Inches(4.3), Inches(3.5), Inches(1.1), RGBColor(0x12,0x4A,0x28))
    add_text(slide, lbl, bx+Inches(0.2), Inches(4.4), Inches(3.1), Inches(0.4),
             size=11, bold=True, color=c, font="Calibri")
    add_text(slide, val, bx+Inches(0.2), Inches(4.8), Inches(3.1), Inches(0.45),
             size=9.5, color=RGBColor(0xBB,0xF7,0xD0), font="Courier New")

# Stack final
add_text(slide,
    "Stack : Next.js 14  ·  Laravel 11  ·  MySQL  ·  Tailwind CSS  ·  Docker  ·  Vercel  ·  Render",
    Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.55),
    size=12, color=RGBColor(0x6B,0x72,0x80), align=PP_ALIGN.CENTER, font="Calibri")

add_text(slide, "11", Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.4),
         size=11, color=GRIS, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SAUVEGARDE
# ══════════════════════════════════════════════════════════════════
out = r"C:\Users\hp\Desktop\BeninMarket_Presentation_Premium.pptx"
prs.save(out)
print(f"✅ Présentation créée : {out}")
print(f"   {len(prs.slides)} slides générées")
